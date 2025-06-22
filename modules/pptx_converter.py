"""
PowerPoint converter module for creating PPTX files from HTML slides.
Provides basic PowerPoint generation functionality for the Streamlit app.
"""

import io
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from modules.models import HTMLSlide, ContentExtractionResult

def create_powerpoint_from_slides(slides: List[HTMLSlide], extraction_result: ContentExtractionResult) -> io.BytesIO:
    """
    Create a PowerPoint presentation from HTML slides.
    
    Args:
        slides: List of HTMLSlide objects
        extraction_result: ContentExtractionResult with document metadata
        
    Returns:
        BytesIO buffer containing the PowerPoint file
    """
    # Create presentation
    prs = Presentation()
    
    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    for i, slide in enumerate(slides):
        # Determine slide layout
        if i == 0:
            # Title slide
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            ppt_slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_placeholder = ppt_slide.shapes.title
            title_placeholder.text = slide.title
            
            # Set subtitle if available
            if len(ppt_slide.placeholders) > 1:
                subtitle_placeholder = ppt_slide.placeholders[1]
                subtitle_placeholder.text = extraction_result.summary
        
        elif i == len(slides) - 1:
            # Conclusion slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            ppt_slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_placeholder = ppt_slide.shapes.title
            title_placeholder.text = slide.title
            
            # Add content
            content_placeholder = ppt_slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.text = "Key Takeaways:"
            
            # Add bullet points (simplified)
            content_lines = slide.section_content.split('\n')
            for line in content_lines[:5]:  # Limit to 5 points
                if line.strip():
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.level = 1
        
        else:
            # Content slide
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            ppt_slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_placeholder = ppt_slide.shapes.title
            title_placeholder.text = slide.title
            
            # Add content
            content_placeholder = ppt_slide.placeholders[1]
            tf = content_placeholder.text_frame
            
            # Extract bullet points from HTML content (simplified)
            content_lines = slide.section_content.split('\n')
            first_line = True
            for line in content_lines[:6]:  # Limit to 6 points
                if line.strip():
                    if first_line:
                        tf.text = line.strip()
                        first_line = False
                    else:
                        p = tf.add_paragraph()
                        p.text = line.strip()
                        p.level = 1
    
    # Apply basic styling
    apply_basic_styling(prs)
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return buffer

def apply_basic_styling(prs: Presentation):
    """Apply basic styling to the presentation."""
    
    # Define color scheme
    primary_color = RGBColor(102, 126, 234)  # Blue
    text_color = RGBColor(44, 62, 80)       # Dark gray
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                
                # Style title shapes
                if shape == slide.shapes.title:
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(36)
                            run.font.color.rgb = primary_color
                            run.font.bold = True
                        paragraph.alignment = PP_ALIGN.CENTER
                
                # Style content shapes
                else:
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                            run.font.color.rgb = text_color
                        paragraph.space_after = Pt(12)

def create_powerpoint_from_template(slides: List[HTMLSlide], extraction_result: ContentExtractionResult, 
                                  template_file: str = None, template_stream: io.BytesIO = None) -> io.BytesIO:
    """
    Create a PowerPoint presentation using a template.
    
    Args:
        slides: List of HTMLSlide objects
        extraction_result: ContentExtractionResult with document metadata
        template_file: Path to template PPTX file
        template_stream: BytesIO stream containing template PPTX
        
    Returns:
        BytesIO buffer containing the PowerPoint file
    """
    # Load template
    if template_stream:
        prs = Presentation(template_stream)
    elif template_file:
        prs = Presentation(template_file)
    else:
        # Fallback to default presentation
        return create_powerpoint_from_slides(slides, extraction_result)
    
    # Clear existing slides (keep layouts)
    slide_layouts = [slide.slide_layout for slide in prs.slides]
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[i]
    
    # Add new slides using template layouts
    for i, slide in enumerate(slides):
        # Choose appropriate layout
        if i == 0 and len(slide_layouts) > 0:
            # Use first layout for title slide
            layout = slide_layouts[0]
        elif len(slide_layouts) > 1:
            # Use second layout for content slides
            layout = slide_layouts[1]
        else:
            # Fallback to first available layout
            layout = slide_layouts[0] if slide_layouts else prs.slide_layouts[0]
        
        ppt_slide = prs.slides.add_slide(layout)
        
        # Populate slide content
        populate_slide_from_template(ppt_slide, slide, i == 0, extraction_result)
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return buffer

def populate_slide_from_template(ppt_slide, html_slide: HTMLSlide, is_title_slide: bool, 
                               extraction_result: ContentExtractionResult):
    """Populate a slide created from template with content."""
    
    # Find title placeholder
    title_shape = None
    content_shape = None
    
    for shape in ppt_slide.shapes:
        if shape.has_text_frame:
            if hasattr(shape, 'placeholder_format'):
                if shape.placeholder_format.type == 1:  # Title placeholder
                    title_shape = shape
                elif shape.placeholder_format.type == 2:  # Content placeholder
                    content_shape = shape
    
    # Set title
    if title_shape:
        title_shape.text = html_slide.title
    
    # Set content
    if content_shape and not is_title_slide:
        tf = content_shape.text_frame
        tf.clear()
        
        # Add content from HTML slide
        content_lines = html_slide.section_content.split('\n')
        first_line = True
        for line in content_lines[:6]:  # Limit to 6 points
            if line.strip():
                if first_line:
                    tf.text = line.strip()
                    first_line = False
                else:
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.level = 1
    elif content_shape and is_title_slide:
        # Add subtitle for title slide
        content_shape.text = extraction_result.summary 