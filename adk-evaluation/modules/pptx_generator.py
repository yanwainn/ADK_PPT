"""
PowerPoint Presentation Generator with Professional Design

This module creates polished .pptx presentations with:
- Professional layouts and themes
- Data visualizations (charts, graphs)
- Corporate branding elements
- Context-aware slide designs
"""

import os
from datetime import datetime
from typing import Dict, Any, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE

import logging
logger = logging.getLogger(__name__)


class PowerPointGenerator:
    """Generate professional PowerPoint presentations with visual elements"""
    
    def __init__(self):
        self.prs = Presentation()
        # Set slide size to 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        # Bosch color palette
        self.colors = {
            'primary': RGBColor(139, 21, 56),      # Bosch red/magenta
            'secondary': RGBColor(0, 169, 206),    # Bosch teal
            'accent': RGBColor(127, 181, 57),      # Bosch green
            'success': RGBColor(127, 181, 57),     # Green
            'text_dark': RGBColor(51, 51, 51),     # Dark gray
            'text_light': RGBColor(255, 255, 255), # White
            'background': RGBColor(255, 255, 255), # White
            'footer_red': RGBColor(237, 28, 36)    # Bosch bright red for footer strip
        }
    
    def add_title_slide(self, title: str, subtitle: str, author: str = ""):
        """Add a Bosch-themed title slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # White background
        left = top = 0
        width = self.prs.slide_width
        height = self.prs.slide_height
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(11.333), Inches(2.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(11.333), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(22)
            subtitle_frame.paragraphs[0].font.color.rgb = self.colors['secondary']  # Blue for subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add Bosch logo text (bottom right)
        logo_box = slide.shapes.add_textbox(
            Inches(10.5), Inches(6.3), Inches(2.5), Inches(0.7)
        )
        logo_frame = logo_box.text_frame
        logo_frame.text = "BOSCH"
        logo_frame.paragraphs[0].font.size = Pt(24)
        logo_frame.paragraphs[0].font.bold = True
        logo_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        logo_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Add colored strip at bottom
        strip_height = Inches(0.15)
        colors = [self.colors['footer_red'], self.colors['secondary'], self.colors['accent']]
        strip_width = width / len(colors)
        
        for i, color in enumerate(colors):
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                strip_width * i, height - strip_height,
                strip_width, strip_height
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = color
            strip.line.fill.background()
    
    def add_content_slide(self, title: str, bullets: List[str], chart_data: Dict = None):
        """Add a content slide with bullets and optional chart"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title (no background, Bosch style)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(11), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']
        
        # Add Bosch logo text (top right)
        logo_box = slide.shapes.add_textbox(
            Inches(11), Inches(0.3), Inches(2), Inches(0.5)
        )
        logo_frame = logo_box.text_frame
        logo_frame.text = "BOSCH"
        logo_frame.paragraphs[0].font.size = Pt(16)
        logo_frame.paragraphs[0].font.bold = True
        logo_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        logo_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Add thin line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3),
            Inches(12.333), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(230, 230, 230)
        line.line.fill.background()
        
        if chart_data:
            # Split layout: bullets on left, chart on right
            # Add bullets with Bosch styling
            bullet_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(6), Inches(5)
            )
            text_frame = bullet_box.text_frame
            text_frame.text = bullets[0] if bullets else ""
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']
            
            for bullet in bullets[1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(16)
                p.font.color.rgb = self.colors['text_dark']
                p.level = 0
            
            # Add chart
            self._add_chart(slide, chart_data, Inches(7), Inches(1.8), Inches(5.8), Inches(4.5))
        else:
            # Full width bullets with Bosch styling
            bullet_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(12.333), Inches(5)
            )
            text_frame = bullet_box.text_frame
            text_frame.text = bullets[0] if bullets else ""
            text_frame.paragraphs[0].font.size = Pt(18)
            text_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']
            text_frame.paragraphs[0].space_before = Pt(12)
            
            for bullet in bullets[1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.font.color.rgb = self.colors['text_dark']
                p.level = 0
                p.space_before = Pt(12)
    
    def add_visual_slide(self, title: str, visual_type: str, data: Dict[str, Any]):
        """Add a slide focused on visual content with Bosch styling"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title (Bosch style)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(11), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']
        
        # Add Bosch logo text (top right)
        logo_box = slide.shapes.add_textbox(
            Inches(11), Inches(0.3), Inches(2), Inches(0.5)
        )
        logo_frame = logo_box.text_frame
        logo_frame.text = "BOSCH"
        logo_frame.paragraphs[0].font.size = Pt(16)
        logo_frame.paragraphs[0].font.bold = True
        logo_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        logo_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Add thin line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3),
            Inches(12.333), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(230, 230, 230)
        line.line.fill.background()
        
        # Add visual based on type
        if visual_type == "chart":
            self._add_chart(slide, data, Inches(1), Inches(1.8), Inches(11.333), Inches(5))
        elif visual_type == "comparison":
            self._add_comparison_table(slide, data, Inches(1), Inches(1.8))
        elif visual_type == "process":
            self._add_process_flow(slide, data, Inches(0.5), Inches(2))
        elif visual_type == "kpi":
            self._add_kpi_cards(slide, data, Inches(0.5), Inches(1.8))
    
    def _add_chart(self, slide, data: Dict, left, top, width, height):
        """Add a chart to the slide"""
        chart_type = data.get('type', 'bar')
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.get('categories', ['Q1', 'Q2', 'Q3', 'Q4'])
        
        # Add series
        series_data = data.get('series', [
            ('Revenue', (65, 78, 82, 91)),
            ('Profit', (55, 65, 70, 80))
        ])
        
        for name, values in series_data:
            chart_data.add_series(name, values)
        
        # Determine chart type
        if chart_type == 'bar':
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type == 'line':
            chart_type_enum = XL_CHART_TYPE.LINE_MARKERS
        elif chart_type == 'pie':
            chart_type_enum = XL_CHART_TYPE.PIE
        else:
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            chart_type_enum, left, top, width, height, chart_data
        ).chart
        
        # Style the chart
        if hasattr(chart, 'has_legend'):
            chart.has_legend = True
            if chart.has_legend and hasattr(chart.legend, 'position'):
                from pptx.enum.chart import XL_LEGEND_POSITION
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    
    def _add_comparison_table(self, slide, data: Dict, left, top):
        """Add a comparison table"""
        rows = len(data.get('rows', [])) + 1  # +1 for header
        cols = len(data.get('columns', []))
        
        if rows == 1 or cols == 0:
            return
        
        # Add table
        width = Inches(11.333)
        height = Inches(0.8 * rows)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Style header row
        for i, col in enumerate(data.get('columns', [])):
            cell = table.cell(0, i)
            cell.text = col
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
            
            # Style text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = self.colors['text_light']
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Add data rows
        for i, row_data in enumerate(data.get('rows', [])):
            for j, value in enumerate(row_data):
                cell = table.cell(i + 1, j)
                cell.text = str(value)
                
                # Alternate row colors
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                
                # Style text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_process_flow(self, slide, data: Dict, left, top):
        """Add a process flow diagram"""
        steps = data.get('steps', ['Research', 'Design', 'Develop', 'Test', 'Deploy'])
        step_width = Inches(2)
        step_height = Inches(1)
        arrow_width = Inches(0.5)
        spacing = Inches(0.2)
        
        current_left = left
        
        for i, step in enumerate(steps):
            # Add step box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                current_left, top + Inches(2), 
                step_width, step_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors['accent']
            
            # Add text
            shape.text = step
            shape.text_frame.paragraphs[0].font.size = Pt(16)
            shape.text_frame.paragraphs[0].font.bold = True
            shape.text_frame.paragraphs[0].font.color.rgb = self.colors['text_light']
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add arrow (except for last step)
            if i < len(steps) - 1:
                arrow_left = current_left + step_width + spacing
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    arrow_left, top + Inches(2.25),
                    arrow_width, Inches(0.5)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.colors['accent']
                
                current_left = arrow_left + arrow_width + spacing
            else:
                current_left += step_width + spacing
    
    def _add_kpi_cards(self, slide, data: Dict, left, top):
        """Add KPI cards"""
        kpis = data.get('kpis', [
            {'value': '$2.5M', 'label': 'Revenue', 'change': '+23%'},
            {'value': '1,234', 'label': 'Customers', 'change': '+15%'},
            {'value': '98.5%', 'label': 'Satisfaction', 'change': '+2%'},
            {'value': '45ms', 'label': 'Response Time', 'change': '-12%'}
        ])
        
        card_width = Inches(2.8)
        card_height = Inches(2)
        spacing = Inches(0.3)
        
        for i, kpi in enumerate(kpis):
            card_left = left + (i % 4) * (card_width + spacing)
            card_top = top + (i // 4) * (card_height + spacing)
            
            # Add card background with Bosch styling
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                card_left, card_top,
                card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(245, 245, 245)
            card.line.color.rgb = self.colors['secondary']
            card.line.width = Pt(2)
            
            # Add value
            value_box = slide.shapes.add_textbox(
                card_left, card_top + Inches(0.3),
                card_width, Inches(0.6)
            )
            value_frame = value_box.text_frame
            value_frame.text = kpi['value']
            value_frame.paragraphs[0].font.size = Pt(28)
            value_frame.paragraphs[0].font.bold = True
            value_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']  # Black for values
            value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add label
            label_box = slide.shapes.add_textbox(
                card_left, card_top + Inches(0.9),
                card_width, Inches(0.4)
            )
            label_frame = label_box.text_frame
            label_frame.text = kpi['label']
            label_frame.paragraphs[0].font.size = Pt(14)
            label_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']
            label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add change
            change_box = slide.shapes.add_textbox(
                card_left, card_top + Inches(1.4),
                card_width, Inches(0.4)
            )
            change_frame = change_box.text_frame
            change_frame.text = kpi['change']
            change_frame.paragraphs[0].font.size = Pt(16)
            change_frame.paragraphs[0].font.bold = True
            color = self.colors['success'] if '+' in kpi['change'] else self.colors['accent']
            change_frame.paragraphs[0].font.color.rgb = color
            change_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def add_conclusion_slide(self, title: str, takeaways: List[str], next_steps: List[str]):
        """Add a conclusion slide with Bosch styling"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(11), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']
        
        # Add Bosch logo text (top right)
        logo_box = slide.shapes.add_textbox(
            Inches(11), Inches(0.3), Inches(2), Inches(0.5)
        )
        logo_frame = logo_box.text_frame
        logo_frame.text = "BOSCH"
        logo_frame.paragraphs[0].font.size = Pt(16)
        logo_frame.paragraphs[0].font.bold = True
        logo_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        logo_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Add thin line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3),
            Inches(12.333), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(230, 230, 230)
        line.line.fill.background()
        
        # Add takeaways section with circular bullet style
        takeaway_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(5.5), Inches(4.5)
        )
        takeaway_frame = takeaway_box.text_frame
        takeaway_frame.text = "Key Takeaways"
        takeaway_frame.paragraphs[0].font.size = Pt(22)
        takeaway_frame.paragraphs[0].font.bold = True
        takeaway_frame.paragraphs[0].font.color.rgb = self.colors['text_dark']  # Black for takeaways
        
        for takeaway in takeaways:
            p = takeaway_frame.add_paragraph()
            p.text = f"• {takeaway}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text_dark']
            p.space_before = Pt(10)
        
        # Add next steps section
        steps_box = slide.shapes.add_textbox(
            Inches(7), Inches(2), Inches(5.5), Inches(4.5)
        )
        steps_frame = steps_box.text_frame
        steps_frame.text = "Next Steps"
        steps_frame.paragraphs[0].font.size = Pt(22)
        steps_frame.paragraphs[0].font.bold = True
        steps_frame.paragraphs[0].font.color.rgb = self.colors['secondary']
        
        for i, step in enumerate(next_steps):
            p = steps_frame.add_paragraph()
            p.text = f"{i+1}. {step}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['text_dark']
            p.space_before = Pt(10)
        
        # Add colored strip at bottom
        strip_height = Inches(0.15)
        colors = [self.colors['footer_red'], self.colors['secondary'], self.colors['accent']]
        strip_width = self.prs.slide_width / len(colors)
        
        for i, color in enumerate(colors):
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                strip_width * i, self.prs.slide_height - strip_height,
                strip_width, strip_height
            )
            strip.fill.solid()
            strip.fill.fore_color.rgb = color
            strip.line.fill.background()
    
    def save(self, filename: str) -> str:
        """Save the presentation and return the path"""
        self.prs.save(filename)
        return os.path.abspath(filename)
    
    def generate_from_content(self, content_data: Dict[str, Any], output_path: str) -> str:
        """Generate complete presentation from structured content"""
        try:
            # Title slide
            title_info = content_data.get('title_slide', {})
            self.add_title_slide(
                title=title_info.get('title', 'Presentation'),
                subtitle=title_info.get('subtitle', ''),
                author=title_info.get('author', '')
            )
            
            # Content slides
            for slide_data in content_data.get('slides', []):
                slide_type = slide_data.get('type', 'content')
                
                if slide_type == 'visual':
                    self.add_visual_slide(
                        title=slide_data.get('title', ''),
                        visual_type=slide_data.get('visual_type', 'chart'),
                        data=slide_data.get('data', {})
                    )
                else:
                    self.add_content_slide(
                        title=slide_data.get('title', ''),
                        bullets=slide_data.get('bullets', []),
                        chart_data=slide_data.get('chart_data')
                    )
            
            # Conclusion slide
            conclusion = content_data.get('conclusion', {})
            if conclusion:
                self.add_conclusion_slide(
                    title=conclusion.get('title', 'Key Takeaways'),
                    takeaways=conclusion.get('takeaways', []),
                    next_steps=conclusion.get('next_steps', [])
                )
            
            # Save presentation
            saved_path = self.save(output_path)
            logger.info(f"✅ PowerPoint presentation saved to: {saved_path}")
            return saved_path
            
        except Exception as e:
            logger.error(f"❌ Error generating PowerPoint: {str(e)}")
            raise